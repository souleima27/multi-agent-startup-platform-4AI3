from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation


@dataclass
class ParsedDocument:
    source_format: str
    extracted_text_preview: str
    extraction_ok: bool
    extraction_notes: list[str]


def _safe_preview(text: str, max_chars: int = 220) -> str:
    compact = " ".join(text.split())
    return compact[:max_chars] if compact else ""


def _read_pdf(path: Path) -> ParsedDocument:
    notes: list[str] = []
    texts: list[str] = []
    try:
        reader = PdfReader(str(path))
        for idx, page in enumerate(reader.pages):
            content = page.extract_text() or ""
            if content.strip():
                texts.append(content)
            else:
                notes.append(f"No extractable text on page {idx + 1}.")
    except Exception as exc:
        return ParsedDocument(
            source_format="pdf",
            extracted_text_preview="",
            extraction_ok=False,
            extraction_notes=[f"PDF extraction failed: {exc}"],
        )

    merged = "\n".join(texts)
    ok = len(merged.strip()) > 0
    if not ok:
        notes.append("PDF appears to be image-based or empty.")
    return ParsedDocument(
        source_format="pdf",
        extracted_text_preview=_safe_preview(merged),
        extraction_ok=ok,
        extraction_notes=notes,
    )


def _read_docx(path: Path) -> ParsedDocument:
    try:
        document = DocxDocument(str(path))
        text = "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text)
    except Exception as exc:
        return ParsedDocument(
            source_format="docx",
            extracted_text_preview="",
            extraction_ok=False,
            extraction_notes=[f"DOCX extraction failed: {exc}"],
        )

    ok = len(text.strip()) > 0
    notes = [] if ok else ["DOCX file contains no readable paragraph text."]
    return ParsedDocument(
        source_format="docx",
        extracted_text_preview=_safe_preview(text),
        extraction_ok=ok,
        extraction_notes=notes,
    )


def _read_pptx(path: Path) -> ParsedDocument:
    texts: list[str] = []
    try:
        deck = Presentation(str(path))
        for slide in deck.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
    except Exception as exc:
        return ParsedDocument(
            source_format="pptx",
            extracted_text_preview="",
            extraction_ok=False,
            extraction_notes=[f"PPTX extraction failed: {exc}"],
        )

    merged = "\n".join(texts)
    ok = len(merged.strip()) > 0
    notes = [] if ok else ["PPTX contains no readable text blocks."]
    return ParsedDocument(
        source_format="pptx",
        extracted_text_preview=_safe_preview(merged),
        extraction_ok=ok,
        extraction_notes=notes,
    )


def parse_document(path: str | Path) -> ParsedDocument:
    doc_path = Path(path)
    suffix = doc_path.suffix.lower()

    if suffix == ".pdf":
        return _read_pdf(doc_path)
    if suffix == ".docx":
        return _read_docx(doc_path)
    if suffix == ".pptx":
        return _read_pptx(doc_path)

    return ParsedDocument(
        source_format="other",
        extracted_text_preview="",
        extraction_ok=False,
        extraction_notes=[f"Unsupported non-image format: {suffix or 'no_extension'}."],
    )
